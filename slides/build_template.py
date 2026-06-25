"""
DANES ScHack 2026 — PowerPoint / Google Slides Template Builder
Run: python3 build_template.py
Produces: ScHack2026_SlideTemplate.pptx  (import into Google Slides via File → Import slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree
import copy, os

# ── Brand colours ────────────────────────────────────────────────
BG       = RGBColor(0xFE, 0xFB, 0xF1)  # #FEFBF1 cream
PRIMARY  = RGBColor(0x52, 0x05, 0x3E)  # #52053E maroon
TEAL     = RGBColor(0x05, 0x87, 0x8A)  # #05878a
YELLOW   = RGBColor(0xFE, 0xDA, 0x75)  # #feda75
DARK     = RGBColor(0x1A, 0x0A, 0x1A)  # near-black
MUTED    = RGBColor(0x6B, 0x52, 0x68)  # muted purple-grey
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
T1       = RGBColor(0x05, 0x87, 0x8A)  # Track 1 teal
T2       = RGBColor(0xB0, 0x60, 0x10)  # Track 2 amber
T3       = RGBColor(0x52, 0x05, 0x3E)  # Track 3 maroon
LIGHT_BG = RGBColor(0xF0, 0xE8, 0xEC)  # light lavender fill

# ── Slide dimensions: 16:9 widescreen ───────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# ── Font names ───────────────────────────────────────────────────
FONT_HEADING = "Georgia"      # closest to Taviraj in standard fonts
FONT_BODY    = "Georgia"
FONT_MONO    = "Courier New"


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ────────────────────────────────────────────────────────────────
# Helper: add a solid-fill rectangle
# ────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE
        if False else 1,  # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE = 5, RECTANGLE = 1
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


# ────────────────────────────────────────────────────────────────
# Helper: add a text box with full control
# ────────────────────────────────────────────────────────────────
def add_text(slide, text, left, top, width, height,
             font_name=FONT_BODY, font_size=18, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


# ────────────────────────────────────────────────────────────────
# Helper: set slide background colour
# ────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ────────────────────────────────────────────────────────────────
# Helper: add a coloured top bar (brand accent)
# ────────────────────────────────────────────────────────────────
def add_top_bar(slide, color, height=Pt(5)):
    bar = add_rect(slide, 0, 0, W, height, color)
    return bar


# ────────────────────────────────────────────────────────────────
# Helper: add footer text (event + slide info)
# ────────────────────────────────────────────────────────────────
def add_footer(slide, left_text="DANES ScHack 2026 · University of Turin",
               right_text="July 6–10, 2026"):
    y = H - Inches(0.32)
    add_text(slide, left_text,  Inches(0.3), y, Inches(7), Inches(0.28),
             font_name=FONT_MONO, font_size=7, color=MUTED)
    add_text(slide, right_text, Inches(6.3), y, Inches(6.7), Inches(0.28),
             font_name=FONT_MONO, font_size=7, color=MUTED, align=PP_ALIGN.RIGHT)
    # thin rule above footer
    rule = add_rect(slide, Inches(0.3), y - Pt(4), W - Inches(0.6), Pt(1),
                    RGBColor(0xD0, 0xB8, 0xCC))


# ────────────────────────────────────────────────────────────────
# Helper: add a filled box with label + body text
# ────────────────────────────────────────────────────────────────
def add_box(slide, label, body, left, top, width, height,
            bg=RGBColor(0xE8, 0xF6, 0xF6), border=TEAL):
    rect = add_rect(slide, left, top, width, height, bg, border)
    if label:
        add_text(slide, label, left + Inches(0.12), top + Inches(0.08),
                 width - Inches(0.24), Inches(0.28),
                 font_size=8, bold=True, color=DARK)
    add_text(slide, body,
             left + Inches(0.12),
             top + (Inches(0.35) if label else Inches(0.1)),
             width - Inches(0.24),
             height - Inches(0.2),
             font_size=13, color=DARK, wrap=True)
    return rect


# ════════════════════════════════════════════════════════════════
# SLIDE 1  Title slide (dark maroon)
# ════════════════════════════════════════════════════════════════
def slide_title(prs, title, subtitle, presenter, track_label="",
                bg_color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, bg_color)

    # decorative yellow bar at very bottom
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), YELLOW)

    # Main title
    add_text(slide, title,
             Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.8),
             font_name=FONT_HEADING, font_size=44, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    if subtitle:
        add_text(slide, subtitle,
                 Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.7),
                 font_size=22, italic=True,
                 color=RGBColor(0xCC, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    # Presenter
    if presenter:
        add_text(slide, presenter,
                 Inches(1.5), Inches(4.55), Inches(10.3), Inches(0.5),
                 font_size=18, color=RGBColor(0xBB, 0xAA, 0xBB),
                 align=PP_ALIGN.CENTER)

    # Meta line
    meta = track_label + ("  ·  " if track_label else "") + "DANES ScHack 2026  ·  University of Turin  ·  July 6–10, 2026"
    add_text(slide, meta,
             Inches(1.5), Inches(5.4), Inches(10.3), Inches(0.4),
             font_name=FONT_MONO, font_size=10,
             color=RGBColor(0x99, 0x80, 0x99), align=PP_ALIGN.CENTER)

    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 2  Section divider
# ════════════════════════════════════════════════════════════════
def slide_section(prs, section_num, section_title, subtitle="",
                  color=TEAL):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, color)

    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), YELLOW)

    if section_num:
        add_text(slide, section_num.upper(),
                 Inches(1.5), Inches(2.2), Inches(10.3), Inches(0.45),
                 font_name=FONT_MONO, font_size=11, bold=True,
                 color=RGBColor(0xCC, 0xEE, 0xEE), align=PP_ALIGN.CENTER)

    add_text(slide, section_title,
             Inches(1.0), Inches(2.8), Inches(11.3), Inches(1.4),
             font_name=FONT_HEADING, font_size=40, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    if subtitle:
        add_text(slide, subtitle,
                 Inches(2.0), Inches(4.35), Inches(9.3), Inches(0.6),
                 font_size=17, italic=True,
                 color=RGBColor(0xDD, 0xF4, 0xF4), align=PP_ALIGN.CENTER)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 3  Standard content — heading + bullets
# ════════════════════════════════════════════════════════════════
def slide_bullets(prs, heading, bullets, note="", bar_color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_top_bar(slide, bar_color)

    add_text(slide, heading,
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.75),
             font_name=FONT_HEADING, font_size=28, bold=True, color=PRIMARY)

    # thin rule under heading
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Pt(1.5),
             RGBColor(0xD0, 0xB8, 0xCC))

    # Bullet list
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.15),
                                     Inches(12.3), Inches(5.2))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after  = Pt(2)
        run = p.add_run()
        # indent sub-bullets (strings starting with "  ")
        if b.startswith("    "):
            p.level = 2
            run.text = b.strip()
            run.font.size = Pt(15)
            run.font.color.rgb = MUTED
        elif b.startswith("  "):
            p.level = 1
            run.text = b.strip()
            run.font.size = Pt(17)
            run.font.color.rgb = DARK
        else:
            p.level = 0
            run.text = b
            run.font.size = Pt(19)
            run.font.color.rgb = DARK
        run.font.name = FONT_BODY

    if note:
        add_box(slide, "Key point", note,
                Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.75),
                bg=RGBColor(0xFF, 0xF5, 0xD6), border=RGBColor(0xC0, 0x90, 0x0A))

    add_footer(slide)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 4  Two-column layout
# ════════════════════════════════════════════════════════════════
def slide_two_col(prs, heading,
                  left_head, left_body,
                  right_head, right_body,
                  bar_color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_top_bar(slide, bar_color)

    add_text(slide, heading,
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.75),
             font_name=FONT_HEADING, font_size=28, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Pt(1.5),
             RGBColor(0xD0, 0xB8, 0xCC))

    # Left column
    add_text(slide, left_head,
             Inches(0.4), Inches(1.1), Inches(5.9), Inches(0.45),
             font_size=17, bold=True, color=TEAL)
    add_text(slide, left_body,
             Inches(0.4), Inches(1.65), Inches(5.9), Inches(5.0),
             font_size=16, color=DARK, wrap=True)

    # Divider
    add_rect(slide, Inches(6.55), Inches(1.1), Pt(1.5), Inches(5.6),
             RGBColor(0xD0, 0xB8, 0xCC))

    # Right column
    add_text(slide, right_head,
             Inches(6.75), Inches(1.1), Inches(6.2), Inches(0.45),
             font_size=17, bold=True, color=TEAL)
    add_text(slide, right_body,
             Inches(6.75), Inches(1.65), Inches(6.2), Inches(5.0),
             font_size=16, color=DARK, wrap=True)

    add_footer(slide)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 5  Code / demo slide
# ════════════════════════════════════════════════════════════════
def slide_code(prs, heading, description, code_text, bar_color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_top_bar(slide, bar_color)

    add_text(slide, heading,
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.75),
             font_name=FONT_HEADING, font_size=28, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Pt(1.5),
             RGBColor(0xD0, 0xB8, 0xCC))

    # Description (left)
    add_text(slide, description,
             Inches(0.4), Inches(1.1), Inches(4.0), Inches(5.5),
             font_size=15, color=DARK, wrap=True)

    # Code block (dark background box)
    code_bg = add_rect(slide,
                       Inches(4.6), Inches(1.1),
                       Inches(8.5), Inches(5.6),
                       RGBColor(0x1E, 0x1E, 0x2E))
    add_text(slide, code_text,
             Inches(4.75), Inches(1.2),
             Inches(8.2), Inches(5.4),
             font_name=FONT_MONO, font_size=12,
             color=RGBColor(0xCC, 0xFF, 0xCC), wrap=True)

    add_footer(slide)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 6  Image placeholder + caption
# ════════════════════════════════════════════════════════════════
def slide_image(prs, heading, caption, source="", bar_color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_top_bar(slide, bar_color)

    add_text(slide, heading,
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.75),
             font_name=FONT_HEADING, font_size=28, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Pt(1.5),
             RGBColor(0xD0, 0xB8, 0xCC))

    # Image placeholder box
    ph = add_rect(slide,
                  Inches(0.4), Inches(1.1),
                  Inches(8.5), Inches(5.4),
                  RGBColor(0xF0, 0xE8, 0xEC),
                  RGBColor(0xCC, 0xA8, 0xC0))
    add_text(slide,
             "[ Replace with image ]\nRecommended: 1200 × 700 px",
             Inches(2.5), Inches(3.2), Inches(4.5), Inches(1.0),
             font_name=FONT_MONO, font_size=12,
             color=RGBColor(0xB0, 0x80, 0xA8), align=PP_ALIGN.CENTER)

    # Caption panel (right)
    add_text(slide, caption,
             Inches(9.2), Inches(1.1), Inches(4.0), Inches(4.5),
             font_size=16, color=DARK, wrap=True)

    if source:
        add_text(slide, f"Source: {source}",
                 Inches(9.2), Inches(5.8), Inches(4.0), Inches(0.4),
                 font_size=10, italic=True, color=MUTED, wrap=True)

    add_footer(slide)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 7  Pull quote
# ════════════════════════════════════════════════════════════════
def slide_quote(prs, heading, quote_text, attribution="", bar_color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_top_bar(slide, bar_color)

    add_text(slide, heading,
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.75),
             font_name=FONT_HEADING, font_size=28, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Pt(1.5),
             RGBColor(0xD0, 0xB8, 0xCC))

    # Yellow left bar
    add_rect(slide, Inches(0.55), Inches(1.25), Inches(0.1), Inches(4.4), YELLOW)

    add_text(slide, quote_text,
             Inches(0.85), Inches(1.3), Inches(11.6), Inches(4.2),
             font_name=FONT_HEADING, font_size=22, italic=True,
             color=PRIMARY, wrap=True)

    if attribution:
        add_text(slide, f"— {attribution}",
                 Inches(0.85), Inches(5.7), Inches(11.6), Inches(0.5),
                 font_size=14, color=MUTED)

    add_footer(slide)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 8  Exercise / hands-on prompt
# ════════════════════════════════════════════════════════════════
def slide_exercise(prs, duration, title, steps, resources="", bar_color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_top_bar(slide, bar_color)

    add_text(slide, f"Exercise  ·  {duration}",
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.75),
             font_name=FONT_HEADING, font_size=28, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Pt(1.5),
             RGBColor(0xD0, 0xB8, 0xCC))

    # Yellow exercise box
    ex_bg = add_rect(slide,
                     Inches(0.4), Inches(1.1),
                     Inches(12.5), Inches(4.8),
                     RGBColor(0xFF, 0xF8, 0xE0),
                     RGBColor(0xFE, 0xDA, 0x75))

    add_text(slide, "⚒  HANDS-ON",
             Inches(0.6), Inches(1.2), Inches(4.0), Inches(0.35),
             font_name=FONT_MONO, font_size=9, bold=True,
             color=RGBColor(0xA0, 0x70, 0x00))

    add_text(slide, title,
             Inches(0.6), Inches(1.6), Inches(12.0), Inches(0.55),
             font_name=FONT_HEADING, font_size=22, bold=True, color=DARK)

    add_text(slide, steps,
             Inches(0.65), Inches(2.3), Inches(8.0), Inches(3.4),
             font_size=16, color=DARK, wrap=True)

    if resources:
        add_box(slide, "Resources", resources,
                Inches(9.0), Inches(2.3), Inches(3.7), Inches(3.4),
                bg=RGBColor(0xE8, 0xF6, 0xF6), border=TEAL)

    add_footer(slide)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 9  Key takeaways
# ════════════════════════════════════════════════════════════════
def slide_takeaways(prs, points, next_up="", bar_color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_top_bar(slide, bar_color)

    add_text(slide, "Key Takeaways",
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.75),
             font_name=FONT_HEADING, font_size=28, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Pt(1.5),
             RGBColor(0xD0, 0xB8, 0xCC))

    for i, pt in enumerate(points):
        y = Inches(1.15) + i * Inches(1.1)
        num_box = add_rect(slide,
                           Inches(0.4), y,
                           Inches(0.55), Inches(0.55),
                           PRIMARY)
        add_text(slide, str(i + 1),
                 Inches(0.4), y + Inches(0.05),
                 Inches(0.55), Inches(0.45),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, pt,
                 Inches(1.1), y, Inches(11.7), Inches(0.95),
                 font_size=17, color=DARK, wrap=True)

    if next_up:
        add_box(slide, "Coming up", next_up,
                Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.72),
                bg=RGBColor(0xFF, 0xF5, 0xD6), border=RGBColor(0xC0, 0x90, 0x0A))

    add_footer(slide)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 10  Hackathon group assignments
# ════════════════════════════════════════════════════════════════
def slide_groups(prs, day, groups, bar_color=PRIMARY):
    """groups = list of dicts: {name, theme, members: [(name, role, lead)]}"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_top_bar(slide, bar_color)

    add_text(slide, f"Hackathon Groups  ·  {day}",
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.75),
             font_name=FONT_HEADING, font_size=28, bold=True, color=PRIMARY)
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Pt(1.5),
             RGBColor(0xD0, 0xB8, 0xCC))

    col_w = Inches(4.0)
    gap   = Inches(0.17)
    for i, g in enumerate(groups[:3]):
        x = Inches(0.4) + i * (col_w + gap)
        # Group header
        hdr = add_rect(slide, x, Inches(1.1), col_w, Inches(0.55), bar_color)
        add_text(slide, g["name"],
                 x + Inches(0.1), Inches(1.13),
                 col_w - Inches(0.2), Inches(0.45),
                 font_size=14, bold=True, color=WHITE)

        add_text(slide, g.get("theme", ""),
                 x + Inches(0.1), Inches(1.75),
                 col_w - Inches(0.2), Inches(0.4),
                 font_size=12, italic=True, color=TEAL)

        # Member list
        member_text = ""
        for m_name, m_role, m_lead in g.get("members", []):
            lead_tag = " [lead]" if m_lead else ""
            member_text += f"• {m_name}{lead_tag}\n  {m_role}\n"

        add_text(slide, member_text.strip(),
                 x + Inches(0.1), Inches(2.25),
                 col_w - Inches(0.2), Inches(4.1),
                 font_size=13, color=DARK, wrap=True)

    add_footer(slide)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 11  Participant presentation — title slide
# ════════════════════════════════════════════════════════════════
def slide_participant_title(prs, project_title, name, institution, track):
    colors = {1: T1, 2: T2, 3: T3}
    bg = colors.get(track, PRIMARY)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, bg)

    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), YELLOW)

    add_text(slide, f"Track {track} Hackathon Presentation",
             Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.45),
             font_name=FONT_MONO, font_size=11,
             color=RGBColor(0xCC, 0xEE, 0xEE), align=PP_ALIGN.CENTER)

    add_text(slide, project_title,
             Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.8),
             font_name=FONT_HEADING, font_size=38, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, name,
             Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.5),
             font_size=20, color=RGBColor(0xEE, 0xDD, 0xEE),
             align=PP_ALIGN.CENTER)

    add_text(slide, institution,
             Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.4),
             font_size=16, italic=True,
             color=RGBColor(0xCC, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    add_text(slide, "DANES ScHack 2026  ·  University of Turin  ·  July 2026",
             Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.4),
             font_name=FONT_MONO, font_size=10,
             color=RGBColor(0x99, 0x80, 0x99), align=PP_ALIGN.CENTER)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 12  Closing / Thank you
# ════════════════════════════════════════════════════════════════
def slide_closing(prs, main_text="Thank You", sub="Questions & Discussion",
                  contact="", bg_color=PRIMARY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, bg_color)

    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), YELLOW)

    add_text(slide, main_text,
             Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.5),
             font_name=FONT_HEADING, font_size=52, bold=True,
             color=YELLOW, align=PP_ALIGN.CENTER)

    if sub:
        add_text(slide, sub,
                 Inches(2.0), Inches(3.95), Inches(9.3), Inches(0.6),
                 font_size=22, italic=True,
                 color=RGBColor(0xDD, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

    if contact:
        add_text(slide, contact,
                 Inches(2.0), Inches(5.1), Inches(9.3), Inches(0.5),
                 font_name=FONT_MONO, font_size=14,
                 color=RGBColor(0xBB, 0xAA, 0xBB), align=PP_ALIGN.CENTER)

    add_text(slide, "opendanes.org  ·  DANES ScHack 2026",
             Inches(2.0), Inches(6.3), Inches(9.3), Inches(0.4),
             font_name=FONT_MONO, font_size=10,
             color=RGBColor(0x88, 0x70, 0x88), align=PP_ALIGN.CENTER)
    return slide


# ════════════════════════════════════════════════════════════════
# BUILD ALL TEMPLATES
# ════════════════════════════════════════════════════════════════

def build_opening():
    prs = new_prs()
    slide_title(prs,
        "DANES ScHack 2026",
        "First DANES Summer School and Hackathon",
        "University of Turin  ·  July 6–10, 2026")

    slide_two_col(prs,
        "What is DANES ScHack?",
        "A summer school",
        "Five days of intensive, hands-on teaching in three parallel tracks — combining morning theory with afternoon hackathon work on your own data.",
        "A hackathon",
        "Each afternoon you work in small groups applying what you learned to real research questions. On Friday you present your findings to the group.",
        bar_color=TEAL)

    slide_bullets(prs,
        "Three Parallel Tracks",
        [
          "Track 1 — Ancient Language Processing  (ALP)",
          "  NLP for cuneiform, Akkadian, Sumerian, Greek, Latin",
          "  Transformer models, embeddings, and fine-tuning",
          "  Prompt engineering, RAG, and chain-of-thought",
          "Track 2 — Computer Vision for Material Culture  (CV)",
          "  Image classification and object detection (YOLO / DETR)",
          "  Handwritten Text Recognition (HTR / OCR)",
          "  3D photogrammetry and surface analysis",
          "Track 3 — Network Analysis for Ancient Cultures  (Networks)",
          "  Social and spatial network analysis",
          "  Dynamic networks: modeling change over time",
          "  Gephi, NetworkX, and interpretation",
        ],
        bar_color=TEAL)

    slide_bullets(prs,
        "The Week at a Glance",
        [
          "Monday 6 July     Opening ceremony · Track introductions",
          "Tuesday 7 July    Morning lectures · Hackathon session 1",
          "Wednesday 8 July  Morning lectures · Hackathon session 2",
          "Thursday 9 July   Morning lectures · Hackathon session 3",
          "Friday 10 July    Final work · Presentations · Closing dinner",
          "",
          "Each day:   09:00–12:30  Morning lectures (all tracks parallel)",
          "            12:30–14:00  Lunch",
          "            14:00–18:00  Hackathon sessions",
          "            19:00        Dinner together",
        ])

    slide_bullets(prs,
        "ScHack 2026 at a Glance",
        [
          "35 confirmed participants",
          "  PhD students, postdoctoral researchers, and faculty",
          "15 countries represented",
          "14+ teachers across three tracks",
          "5 days · 3 tracks · 4 hackathon sessions",
          "",
          "This is the first edition — your feedback will shape future years.",
        ],
        note="No prior programming experience required. Participants are encouraged to bring their own datasets and research questions.")

    slide_two_col(prs,
        "Practical Information",
        "Computing",
        "• Bring your own laptop\n• API access (OpenAI / Anthropic) provided — details on Day 1\n• Google Colab for GPU sessions — you need a Google account\n• Wi-Fi: [NETWORK]  Password: [PASSWORD]\n\nCommunication\n• We will use [Slack / Discord] — join link on the table\n• Schedule and materials: opendanes.org",
        "Meals & Accommodation",
        "• Breakfast and dinner at the residence\n• Lunch provided on-site\n• Vegetarian options at all meals\n• Flag dietary needs today if not already done\n\nOrganisers\n• Maurizio Viano — maurizio.viano@unito.it\n• Elena Devecchi — elena.devecchi@unito.it\n• Shai Gordin — shygordin@gmail.com",
        bar_color=TEAL)

    slide_closing(prs,
        "Let's begin.",
        "Welcome to Turin. Welcome to DANES ScHack 2026.",
        "opendanes.org")

    out = "ScHack2026_Opening.pptx"
    prs.save(out)
    print(f"  ✓  {out}")
    return out


def build_lecture():
    prs = new_prs()

    slide_title(prs,
        "[Session Title]",
        "[One-line description of what this session covers]",
        "[Your Name]  ·  [Institution]",
        track_label="Track 1 — Ancient Language Processing",
        bg_color=T1)

    slide_bullets(prs,
        "Session Outline",
        [
          "Part 1   [First topic]                   ~20 min",
          "Part 2   [Second topic]                  ~25 min",
          "Part 3   [Third topic]                   ~20 min",
          "Exercise  Hands-on practice              ~20 min",
          "Q & A    Discussion",
          "",
          "After this session you will be able to:",
          "  [Learning outcome 1]",
          "  [Learning outcome 2]",
          "  [Learning outcome 3]",
        ],
        note="Connection to hackathon: [How today's content feeds into the afternoon session]",
        bar_color=T1)

    slide_section(prs,
        "Part 1", "[Section Title]",
        "[One sentence describing what this section covers]",
        color=T1)

    slide_two_col(prs,
        "[Key Concept Name]",
        "Definition",
        "[A clear, precise definition in 2–3 sentences. Use plain language before technical terms. Explain what problem this concept solves and why it arose.]",
        "Also known as / Related",
        "• [Synonym or variant term]\n• [Related concept]\n\nNot to be confused with:\n• [Common misconception]\n\nSee also:\n• [Relevant paper or resource]",
        bar_color=T1)

    slide_bullets(prs,
        "[Slide Title — Main Content]",
        [
          "[Point 1 — the most important]",
          "  [Sub-point or clarification]",
          "[Point 2]",
          "[Point 3]",
          "[Point 4]",
          "[Point 5 — max 5 bullets per slide]",
        ],
        note="Key takeaway: [The one thing participants must remember from this slide]",
        bar_color=T1)

    slide_two_col(prs,
        "[Comparison Title — e.g. Fine-tuning vs. Prompting]",
        "[Left column heading]",
        "• [Point A1]\n• [Point A2]\n• [Point A3]\n\nUse when: [When to prefer this approach]",
        "[Right column heading]",
        "• [Point B1]\n• [Point B2]\n• [Point B3]\n\nUse when: [When to prefer this approach]",
        bar_color=T1)

    slide_code(prs,
        "[Code / Demo Title]",
        "[Describe what this code demonstrates. Reference an ancient corpus or dataset. Note 1–2 things to watch when running it.]",
        """# [Replace with your actual code]
from anthropic import Anthropic

client = Anthropic()

message = client.messages.create(
    model="claude-sonnet-4-6",
    max_tokens=1024,
    messages=[{
        "role": "user",
        "content": "[Your prompt here]"
    }]
)
print(message.content[0].text)""",
        bar_color=T1)

    slide_image(prs,
        "[Slide Title]",
        "[Explanation of what the image shows and why it matters for the session. Point out 1–2 specific features participants should notice.]",
        source="[Image credit / licence]",
        bar_color=T1)

    slide_quote(prs,
        "[Source or Context — e.g. 'ARM 26 1, ll. 1–5']",
        "[Primary source text, transliteration, or expert quotation]",
        attribution="[Author, text, line numbers, date, collection]",
        bar_color=T1)

    slide_exercise(prs,
        duration="20 minutes",
        title="[Exercise title / prompt]",
        steps="1. [Open / load — notebook or dataset]\n\n2. [Run / modify — what to change or try]\n\n3. [Observe / compare — what output to look at]\n\n4. [Optional stretch goal]",
        resources="Notebook:\n[filename or Colab link]\n\nData:\n[corpus / dataset]\n\nDocs:\n[relevant URL]",
        bar_color=T1)

    slide_takeaways(prs,
        points=[
          "[Takeaway 1] — [brief elaboration of why this matters]",
          "[Takeaway 2] — [brief elaboration]",
          "[Takeaway 3] — [brief elaboration]",
        ],
        next_up="[What the next session / hackathon will build on from today]")

    slide_closing(prs,
        "Questions?",
        "[Optional: a discussion question to open the floor]",
        "[Your Name]  ·  [your@email.com]",
        bg_color=T1)

    out = "ScHack2026_LectureTemplate.pptx"
    prs.save(out)
    print(f"  ✓  {out}")
    return out


def build_hackathon():
    prs = new_prs()

    slide_title(prs,
        "Hackathon Session [N]",
        "[Day / Date]  ·  14:00 – 18:00",
        "DANES ScHack 2026",
        bg_color=PRIMARY)

    slide_bullets(prs,
        "Goals for Today",
        [
          "[Goal 1 — the main skill or method you will practise]",
          "[Goal 2 — the dataset or corpus you will apply it to]",
          "[Goal 3 — the output or result you aim to produce]",
          "",
          "By 18:00 each group should have:",
          "  A working notebook or script demonstrating the approach",
          "  A brief write-up (bullet points) of what worked and what didn't",
          "  One result to share at tomorrow's opening",
        ],
        note="It is fine — and scientifically valuable — to document what didn't work and why.")

    slide_groups(prs,
        day="Session [N]",
        groups=[
          {
            "name": "Group A — [Theme]",
            "theme": "[Method focus]",
            "members": [
              ("[Name 1]", "[Level · Institution]", True),
              ("[Name 2]", "[Level · Institution]", False),
              ("[Name 3]", "[Level · Institution]", False),
              ("[Name 4]", "[Level · Institution]", False),
            ]
          },
          {
            "name": "Group B — [Theme]",
            "theme": "[Method focus]",
            "members": [
              ("[Name 1]", "[Level · Institution]", True),
              ("[Name 2]", "[Level · Institution]", False),
              ("[Name 3]", "[Level · Institution]", False),
            ]
          },
          {
            "name": "Group C — [Theme]",
            "theme": "[Method focus]",
            "members": [
              ("[Name 1]", "[Level · Institution]", True),
              ("[Name 2]", "[Level · Institution]", False),
              ("[Name 3]", "[Level · Institution]", False),
            ]
          },
        ])

    slide_two_col(prs,
        "Task Brief",
        "The challenge\n\n[Describe the specific task or research question for this session in 3–5 sentences. Be concrete — participants should know exactly what they are trying to produce.]\n\nAnchor dataset (if no own data)\n\n[Name of the shared dataset, where to find it, what it contains]",
        "Suggested approach\n\n1. [Step 1 — data loading / exploration]\n\n2. [Step 2 — apply the method]\n\n3. [Step 3 — evaluate / interpret]\n\nStretch goal\n\n[Optional harder task for groups that finish early]")

    slide_bullets(prs,
        "Datasets & Tools",
        [
          "Shared datasets available at: [URL / shared folder]",
          "  [Dataset 1] — [description, language, size]",
          "  [Dataset 2] — [description]",
          "  [Dataset 3] — [description]",
          "",
          "Tools and access",
          "  API key (LiteLLM proxy): [endpoint URL]  — see handout",
          "  Google Colab notebooks: [shared folder link]",
          "  Documentation: [relevant docs URL]",
          "",
          "If your own data is not ready, use [anchor dataset]",
        ],
        note="API spending limit per group: [e.g. $10]. Be mindful — use GPT-4o-mini for iteration, GPT-4o for final runs.")

    slide_bullets(prs,
        "Timeline",
        [
          "14:00–14:20   Task brief and group formation",
          "14:20–16:00   Main work sprint — apply the method",
          "16:00–16:15   Coffee break",
          "16:15–17:30   Continued work / refinement",
          "17:30–18:00   Write-up — document your results (bullet points)",
          "",
          "Deliverable by 18:00:",
          "  Notebook saved and shared to the group folder",
          "  2–3 bullet-point summary of findings",
        ],
        note="Tomorrow's opening: one person per group gives a 3-minute verbal update on what you found.")

    slide_closing(prs,
        "Go!",
        "Find your group · Open your notebook · Start exploring.",
        bg_color=T1)

    out = "ScHack2026_HackathonBrief.pptx"
    prs.save(out)
    print(f"  ✓  {out}")
    return out


def build_participant():
    prs = new_prs()

    slide_participant_title(prs,
        "[Your Project Title]",
        "[Your Name]",
        "[Your Institution]",
        track=1)  # change to 2 or 3

    slide_bullets(prs,
        "Research Question",
        [
          "[State your core research question in one or two sentences.]",
          "",
          "Dataset",
          "  [Name and brief description of the corpus or dataset you worked with]",
          "  [Scale: how many texts / objects / records?]",
          "  [Source: ORACC / CDLI / museum collection / your own data]",
          "",
          "Why this question?",
          "  [1–2 sentences on why this matters for your research or the field]",
        ],
        bar_color=T1)

    slide_two_col(prs,
        "Methods & Approach",
        "What I tried\n\n1. [Method or tool — first approach]\n\n2. [Method or tool — second approach or step]\n\n3. [Any comparison or iteration]",
        "Key decisions\n\n• [Why did you choose this method?]\n\n• [What parameters or settings did you use?]\n\n• [What did you leave out and why?]",
        bar_color=T1)

    slide_two_col(prs,
        "Results",
        "What worked\n\n[Describe what the method produced — numbers, examples, patterns, visualisations. Be specific.]\n\nExample output:\n\n[Paste a short example result here]",
        "What didn't work\n\n[Honest account of failures, errors, or unexpected outputs — this is just as valuable]\n\nLimitations\n\n[What would you need to address to get better results?]",
        bar_color=T1)

    slide_quote(prs,
        "Key Finding",
        "[State your most interesting finding in one or two sentences. It can be a positive result, a negative result, or a methodological insight.]",
        attribution="[Your Name], DANES ScHack 2026",
        bar_color=T1)

    slide_bullets(prs,
        "Discussion",
        [
          "[What does this result mean for your research question?]",
          "[How does it connect to the broader method or field?]",
          "[What surprised you?]",
          "",
          "What I learned this week",
          "  [The most important new skill or concept you took away]",
          "  [Something that changed how you think about the problem]",
        ],
        bar_color=T1)

    slide_bullets(prs,
        "Next Steps",
        [
          "[What would you do differently if you had more time?]",
          "[What is the next concrete step in your research?]",
          "[Are there collaborations or resources from this week you plan to follow up on?]",
        ],
        note="If you would like to continue this work with others from the school, note that here — this is one of the key purposes of ScHack.",
        bar_color=T1)

    slide_closing(prs,
        "Thank You",
        "Questions & Discussion",
        "[your@email.com]",
        bg_color=T1)

    out = "ScHack2026_ParticipantPresentation.pptx"
    prs.save(out)
    print(f"  ✓  {out}")
    return out


# ── Main ─────────────────────────────────────────────────────────
if __name__ == "__main__":
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    print("Building DANES ScHack 2026 slide templates...\n")
    build_opening()
    build_lecture()
    build_hackathon()
    build_participant()
    print("\nDone. Import any .pptx into Google Slides via File → Import slides.")
